from pptx import Presentation
import sys

def extract_text(ppt_path):
    try:
        prs = Presentation(ppt_path)
        for i, slide in enumerate(prs.slides):
            print(f"--- Slide {i+1} ---")
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    print(shape.text)
    except Exception as e:
        print(f"Error: {e}")

if __name__ == "__main__":
    extract_text("aicomputervision (1).pptx")
